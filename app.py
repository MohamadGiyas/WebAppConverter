import os
import sys
import zipfile
import subprocess
import shutil
import uuid
import threading
import time

# Flask and Werkzeug
from flask import (Flask, render_template, request, send_file, 
                   flash, redirect, url_for, jsonify)
from werkzeug.utils import secure_filename

# PDF and Document Libraries
import PyPDF2
import fitz  # PyMuPDF
import pdfplumber
from pdf2docx import Converter

# Image Library
from PIL import Image

# Spreadsheet Library
import pandas as pd

# ReportLab for PDF Creation
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import cm
from reportlab.lib.colors import Color as ReportLabColor

# PowerPoint Libraries
from pptx import Presentation
from pptx.util import Inches

# Ghostscript Library (pastikan di-import, meskipun tidak dipanggil langsung)

# ==============================================================================
# 1. KONFIGURASI APLIKASI
# ==============================================================================
app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get(
    'SECRET_KEY',
    'kunci-rahasia-super-aman-giyas'
)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def allowed_file(filename, allowed_extensions):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in allowed_extensions


def cleanup_stale_files():
    """
    Membersihkan file sisa dari proses sebelumnya.
    File yang sedang dikunci akan dilewati.
    """
    for folder in [UPLOAD_FOLDER]:
        if not os.path.exists(folder):
            continue

        for filename in os.listdir(folder):
            path = os.path.join(folder, filename)

            if not os.path.isfile(path):
                continue

            try:
                os.remove(path)
            except (PermissionError, OSError):
                # File yang masih digunakan tidak dipaksakan untuk dihapus.
                pass


def send_file_with_cleanup(file_path, download_name, cleanup_paths):
    """
    Mengirim file lalu membersihkannya setelah response selesai.

    Pada Windows, beberapa library seperti pdf2docx/PyMuPDF dan file
    response Flask dapat melepas file handle sedikit setelah callback
    response dipanggil. Karena itu cleanup dijalankan di background
    dengan delay awal, lalu retry beberapa kali.
    """

    response = send_file(
        file_path,
        as_attachment=True,
        download_name=download_name
    )

    def cleanup():
        # Beri waktu agar file handle milik Flask/browser/library benar-benar
        # dilepas sebelum Windows mencoba menghapus file.
        time.sleep(2.0)

        paths = [path for path in cleanup_paths if path]

        for attempt in range(20):
            remaining = []

            for path in paths:
                if not os.path.exists(path):
                    continue

                try:
                    os.remove(path)
                    print(f"File cleanup berhasil: {path}")

                except PermissionError:
                    remaining.append(path)
                    print(
                        f"File masih digunakan, percobaan "
                        f"{attempt + 1}/20: {path}"
                    )

                except OSError as error:
                    if getattr(error, "winerror", None) == 32:
                        remaining.append(path)
                        print(
                            f"File masih terkunci, percobaan "
                            f"{attempt + 1}/20: {path}"
                        )
                    else:
                        print(
                            f"Gagal menghapus file {path}: {error}"
                        )

            if not remaining:
                return

            paths = remaining
            time.sleep(0.5)

        for path in paths:
            if os.path.exists(path):
                print(
                    f"File belum dapat dihapus setelah 20 percobaan: {path}"
                )

    def start_cleanup():
        # Jalankan cleanup di thread terpisah supaya proses pengiriman
        # response tidak tertahan oleh retry cleanup.
        thread = threading.Thread(
            target=cleanup,
            daemon=True
        )
        thread.start()

    response.call_on_close(start_cleanup)
    return response


# ==============================================================================
# 2. ROUTE UTAMA
# ==============================================================================
@app.route('/')
def index():
    return render_template('index.html')

# ==============================================================================
# 3. FITUR-FITUR KONVERSI
# ==============================================================================

# Fitur 3.1: Konverter File Umum (Word/PDF/JPG)

# Ganti hanya fungsi /convert di app.py Anda

@app.route('/convert', methods=['POST'])
def convert():
    file = request.files.get('file')
    conversion_type = request.form.get('conversion')

    if not file or file.filename == '':
        return jsonify({'error': 'Tidak ada file yang dipilih!'}), 400

    filename = secure_filename(file.filename)
    
    if conversion_type == 'word-to-pdf' and filename.lower().endswith('.doc'):
        return jsonify({'error': 'Format .doc tidak didukung. Harap simpan sebagai .docx terlebih dahulu.'}), 400

    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_path = ""

    try:
        if conversion_type == 'word-to-pdf':
            if sys.platform != 'win32':
                raise RuntimeError(
                    'Fitur Word ke PDF pada versi lokal ini membutuhkan '
                    'Microsoft Word di Windows.'
                )

            import comtypes
            from docx2pdf import convert as convert_word_to_pdf

            comtypes.CoInitialize()
            try:
                output_filename = filename.replace('.docx', '.pdf')
                output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
                convert_word_to_pdf(input_path, output_path)
            finally:
                comtypes.CoUninitialize()

        elif conversion_type == 'pdf-to-word':
            output_filename = filename.replace('.pdf', '.docx')
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()

        elif conversion_type == 'jpg-to-pdf':
            output_filename = os.path.splitext(filename)[0] + '.pdf'
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
            img = Image.open(input_path).convert("RGB")
            img.save(output_path)
        
        else:
            return jsonify({'error': 'Tipe konversi tidak valid.'}), 400

        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        if os.path.exists(input_path):
             os.remove(input_path)
        print(f"Error during conversion: {e}")
        return jsonify({'error': f'Terjadi kesalahan saat konversi: {str(e)}'}), 500

# Fitur 3.2: Gabung PDF
@app.route('/merge_pdf', methods=['POST'])
def merge_pdf():
    files = request.files.getlist('files[]')

    if len(files) < 2:
        flash('Pilih minimal 2 file untuk digabungkan.', 'danger')
        return redirect(url_for('index'))

    input_paths = []
    temp_pdf_paths = []
    output_path = ""

    try:
        merger = PyPDF2.PdfMerger()

        for file in files:

            if not file or file.filename == "":
                continue

            filename = secure_filename(file.filename)
            extension = filename.rsplit('.', 1)[1].lower()

            input_path = os.path.join(
                app.config['UPLOAD_FOLDER'],
                f"{uuid.uuid4()}_{filename}"
            )

            file.save(input_path)
            input_paths.append(input_path)

            # ===========================
            # Jika file PDF
            # ===========================
            if extension == "pdf":
                merger.append(input_path)

            # ===========================
            # Jika file gambar
            # ===========================
            elif extension in ["jpg", "jpeg", "png", "bmp", "webp"]:

                img = Image.open(input_path)

                if img.mode != "RGB":
                    img = img.convert("RGB")

                temp_pdf = os.path.join(
                    app.config['UPLOAD_FOLDER'],
                    f"{uuid.uuid4()}.pdf"
                )

                img.save(temp_pdf)

                temp_pdf_paths.append(temp_pdf)

                merger.append(temp_pdf)

            else:
                continue

        output_filename = "merged_document.pdf"

        output_path = os.path.join(
            app.config['UPLOAD_FOLDER'],
            output_filename
        )

        merger.write(output_path)
        merger.close()

        return send_file_with_cleanup(
            output_path,
            output_filename,
            input_paths + temp_pdf_paths + [output_path]
        )

    except Exception as e:

        for path in input_paths + temp_pdf_paths:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except:
                pass

        flash(f"Terjadi kesalahan saat menggabungkan file: {e}", "danger")
        return redirect(url_for('index'))
    
# Fitur 3.3: Pisah PDF
@app.route('/split_pdf', methods=['POST'])
def split_pdf():
    file = request.files.get('file')
    pages_str = request.form.get('pages')
    if not file or file.filename == '':
        flash('Tidak ada file yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_path = ""

    try:
        reader = PyPDF2.PdfReader(input_path)
        writer = PyPDF2.PdfWriter()
        page_numbers = set()
        for part in pages_str.replace(' ', '').split(','):
            if '-' in part:
                start, end = map(int, part.split('-'))
                for i in range(start, end + 1):
                    if 1 <= i <= len(reader.pages): page_numbers.add(i - 1)
            else:
                page = int(part)
                if 1 <= page <= len(reader.pages): page_numbers.add(page - 1)
        
        if not page_numbers:
            flash('Tidak ada halaman yang valid untuk dipisah.', 'danger')
            return redirect(url_for('index'))

        for page_num in sorted(list(page_numbers)):
            writer.add_page(reader.pages[page_num])

        output_filename = f"{os.path.splitext(filename)[0]}_split.pdf"
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat memisah PDF: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.4: Tambah Watermark
@app.route('/add_watermark', methods=['POST'])
def add_watermark():
    file = request.files.get('file')
    watermark_text = request.form.get('watermark_text', 'Mohamad Giyas')
    if not file or file.filename == '':
        flash('Tidak ada file yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    temp_watermark_pdf = os.path.join(app.config['UPLOAD_FOLDER'], f"watermark_{os.getpid()}.pdf")
    output_path = ""
    
    try:
        c = canvas.Canvas(temp_watermark_pdf, pagesize=letter)
        page_width, page_height = letter
        c.setFont('Helvetica-Bold', 40)
        c.setFillColor(ReportLabColor(0.8, 0.8, 0.8, alpha=0.5))
        c.translate(page_width/2, page_height/2)
        c.rotate(45)
        text_width = c.stringWidth(watermark_text, 'Helvetica-Bold', 40)
        c.drawString(-text_width/2, -20, watermark_text)
        c.save()

        reader = PyPDF2.PdfReader(input_path)
        watermark_reader = PyPDF2.PdfReader(temp_watermark_pdf)
        watermark_page = watermark_reader.pages[0]
        writer = PyPDF2.PdfWriter()

        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)

        output_filename = f"{os.path.splitext(filename)[0]}_watermarked.pdf"
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, temp_watermark_pdf, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat menambahkan watermark: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.5: PDF ke JPG
@app.route('/pdf_to_jpg', methods=['POST'])
def pdf_to_jpg():
    file = request.files.get('file')
    if not file or file.filename == '':
        flash('Tidak ada file yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_zip_path = ""
    image_paths = []

    try:
        doc = fitz.open(input_path)
        base_filename = os.path.splitext(filename)[0]
        for i, page in enumerate(doc):
            pix = page.get_pixmap()
            image_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{base_filename}_page_{i+1}.jpg")
            pix.save(image_path)
            image_paths.append(image_path)
        doc.close()

        output_zip_filename = f"{base_filename}_images.zip"
        output_zip_path = os.path.join(app.config['UPLOAD_FOLDER'], output_zip_filename)
        with zipfile.ZipFile(output_zip_path, 'w') as zipf:
            for image_path in image_paths:
                zipf.write(image_path, os.path.basename(image_path))
        
        return send_file_with_cleanup(
            output_zip_path,
            output_zip_filename,
            [input_path, output_zip_path] + image_paths
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat konversi PDF ke JPG: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.6: Excel ke PDF
@app.route('/excel_to_pdf', methods=['POST'])
def excel_to_pdf():
    file = request.files.get('file')
    if not file or file.filename == '':
        flash('Tidak ada file Excel yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_path = ""

    try:
        df = pd.read_excel(input_path)
        output_filename = f"{os.path.splitext(filename)[0]}.pdf"
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
        
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        data = [df.columns.tolist()] + df.values.tolist()
        table = Table(data)
        style = TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.grey),
            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('GRID', (0,0), (-1,-1), 1, colors.black)
        ])
        table.setStyle(style)
        doc.build([table])
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat konversi Excel ke PDF: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.7: PDF ke Excel
@app.route('/pdf_to_excel', methods=['POST'])
def pdf_to_excel():
    file = request.files.get('file')
    if not file or file.filename == '':
        flash('Tidak ada file PDF yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_path = ""

    try:
        all_tables = []
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    all_tables.append(pd.DataFrame(table[1:], columns=table[0]))
        
        if not all_tables:
            flash('Tidak ada tabel yang bisa dideteksi di dalam file PDF ini.', 'warning')
            return redirect(url_for('index'))

        output_filename = f"{os.path.splitext(filename)[0]}.xlsx"
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for i, df in enumerate(all_tables):
                df.to_excel(writer, sheet_name=f'Tabel_{i+1}', index=False)

        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat konversi PDF ke Excel: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.8: Konversi Spreadsheet (Excel <-> CSV)
@app.route('/convert_spreadsheet', methods=['POST'])
def convert_spreadsheet():
    file = request.files.get('file')
    conversion_type = request.form.get('conversion')
    if not file or file.filename == '':
        flash('Tidak ada file yang dipilih!', 'danger')
        return redirect(url_for('index'))
    
    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_path = ""

    try:
        if conversion_type == 'excel-to-csv':
            df = pd.read_excel(input_path)
            output_filename = f"{os.path.splitext(filename)[0]}.csv"
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
            df.to_csv(output_path, index=False)
        elif conversion_type == 'csv-to-excel':
            df = pd.read_csv(input_path)
            output_filename = f"{os.path.splitext(filename)[0]}.xlsx"
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
            df.to_excel(output_path, index=False, engine='openpyxl')
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat konversi: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.9: Kompres PDF
@app.route('/compress_pdf', methods=['POST'])
def compress_pdf():
    file = request.files.get('file')
    if not file or file.filename == '':
        flash('Tidak ada file PDF yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_filename = f"{os.path.splitext(filename)[0]}_compressed.pdf"
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)

    try:
        gs_executable = shutil.which(
            'gswin64c' if sys.platform == 'win32' else 'gs'
        )

        if not gs_executable:
            flash(
                'Ghostscript tidak tersedia di environment ini. '
                'Fitur kompres PDF membutuhkan Ghostscript.',
                'danger'
            )
            return redirect(url_for('index'))

        gs_command = [
            gs_executable,
            '-sDEVICE=pdfwrite',
            '-dCompatibilityLevel=1.4',
            '-dPDFSETTINGS=/ebook',
            '-dNOPAUSE',
            '-dQUIET',
            '-dBATCH',
            f'-sOutputFile={output_path}',
            input_path
        ]
        subprocess.run(gs_command, check=True)
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except subprocess.CalledProcessError:
        flash('Gagal menjalankan Ghostscript. Pastikan sudah terinstal dan ditambahkan ke PATH sistem.', 'danger')
        return redirect(url_for('index'))
    except Exception as e:
        flash(f"Terjadi kesalahan saat kompresi PDF: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.10: Kompres Foto
@app.route('/compress_image', methods=['POST'])
def compress_image():
    file = request.files.get('file')
    quality = int(request.form.get('quality', 85))
    resize_percent = int(request.form.get('resize_percent', 100))

    if not file or file.filename == '':
        flash('Tidak ada file gambar yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    base_filename, ext = os.path.splitext(filename)
    output_filename = f"{base_filename}_compressed{ext}"
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)

    try:
        img = Image.open(input_path)
        if resize_percent < 100:
            w, h = img.size
            new_w = int(w * (resize_percent / 100))
            new_h = int(h * (resize_percent / 100))
            img = img.resize((new_w, new_h), Image.LANCZOS)

        if ext.lower() in ['.jpg', '.jpeg']:
            img.save(output_path, optimize=True, quality=quality)
        else:
            img.save(output_path, optimize=True)
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat kompresi foto: {e}", 'danger')
        return redirect(url_for('index'))

# Fitur 3.11: PowerPoint ke PDF
@app.route('/pptx_to_pdf', methods=['POST'])
def pptx_to_pdf():
    if sys.platform != 'win32':
        flash('Fitur ini hanya berfungsi di Windows dengan PowerPoint terinstall.', 'danger')
        return redirect(url_for('index'))

    file = request.files.get('file')
    if not file or file.filename == '':
        flash('Tidak ada file PowerPoint yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_filename = f"{os.path.splitext(filename)[0]}.pdf"
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
    powerpoint = None
    
    try:
        import comtypes
        import comtypes.client

        comtypes.CoInitialize()
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        deck = powerpoint.Presentations.Open(os.path.abspath(input_path))
        deck.SaveAs(os.path.abspath(output_path), 32)
        deck.Close()
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path]
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat konversi. Pastikan PowerPoint terinstall. Error: {e}", 'danger')
        return redirect(url_for('index'))
    finally:
        if powerpoint:
            try:
                powerpoint.Quit()
            except Exception:
                pass

        try:
            comtypes.CoUninitialize()
        except Exception:
            pass

# Fitur 3.12: PDF ke PowerPoint
@app.route('/pdf_to_pptx', methods=['POST'])
def pdf_to_pptx():
    file = request.files.get('file')
    if not file or file.filename == '':
        flash('Tidak ada file PDF yang dipilih!', 'danger')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(input_path)
    output_filename = f"{os.path.splitext(filename)[0]}.pptx"
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
    image_paths = []

    try:
        doc = fitz.open(input_path)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=300)
            image_path = os.path.join(app.config['UPLOAD_FOLDER'], f"slide_{i}.png")
            pix.save(image_path)
            image_paths.append(image_path)
        doc.close()

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6]

        for image_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        prs.save(output_path)
        
        return send_file_with_cleanup(
            output_path,
            output_filename,
            [input_path, output_path] + image_paths
        )

    except Exception as e:
        flash(f"Terjadi kesalahan saat konversi PDF ke PPTX: {e}", 'danger')
        return redirect(url_for('index'))

# ==============================================================================
# ERROR HANDLER
# ==============================================================================

@app.errorhandler(413)
def request_entity_too_large(error):
    return jsonify({
        'error': 'Ukuran file terlalu besar. Maksimal 50 MB.'
    }), 413


# ==============================================================================
# 4. MENJALANKAN APLIKASI
# ==============================================================================
if __name__ == '__main__':
    cleanup_stale_files()
    port = int(os.environ.get('PORT', 5000))
    app.run(
        host='0.0.0.0',
        port=port,
        debug=False
    )